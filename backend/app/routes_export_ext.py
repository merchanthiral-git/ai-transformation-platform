"""Extended export routes — docx generation, templates, module data export."""

from io import BytesIO
import pandas as pd
from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import StreamingResponse

from app.store import store
from app.helpers import get_series, safe_value_counts, dataframe_to_excel_bytes
from app.shared import _safe, _j, _f

router = APIRouter(prefix="/api", tags=["export-ext"])

@router.get("/export/template/{template_type}")
async def download_template(template_type: str):
    """Download Excel template with correct column headers."""
    import io
    from openpyxl import Workbook
    from fastapi.responses import StreamingResponse
    
    wb = Workbook()
    ws = wb.active
    
    templates = {
        "workforce": {
            "sheet_name": "Workforce",
            "headers": ["Employee ID", "Name", "Job Title", "Function", "Sub-Function", "Career Level", "Career Track", "Manager ID", "Compensation", "Tenure", "Location"],
            "samples": [
                ["EMP001", "Sarah Chen", "Data Analyst", "Technology", "Analytics", "L3", "IC", "EMP010", 85000, 3, "New York"],
                ["EMP002", "James Rodriguez", "Senior Data Analyst", "Technology", "Analytics", "L4", "IC", "EMP010", 105000, 5, "New York"],
                ["EMP003", "Maria Thompson", "Financial Analyst", "Finance", "FP&A", "L3", "IC", "EMP011", 90000, 4, "Chicago"],
                ["EMP004", "David Kim", "Operations Coordinator", "Operations", "Supply Chain", "L2", "IC", "EMP012", 65000, 2, "Chicago"],
                ["EMP010", "Rachel Anderson", "VP Technology", "Technology", "Leadership", "L5", "Manager", "EMP020", 185000, 8, "San Francisco"],
            ],
            "notes": "Required for: Workforce Snapshot, Job Architecture, Headcount Planning. Career Level: L1-L6. Career Track: IC or Manager.",
        },
        "work_design": {
            "sheet_name": "Work Design",
            "headers": ["Job Title", "Task Name", "Est Hours/Week", "Current Time Spent %", "Task Type", "Logic", "Interaction", "AI Impact", "Primary Skill", "Secondary Skill"],
            "samples": [
                ["Data Analyst", "Data extraction & cleaning", 10, 25, "Repetitive", "Deterministic", "Independent", "High", "Data Analysis", "Python/SQL"],
                ["Data Analyst", "Report generation & formatting", 8, 20, "Repetitive", "Deterministic", "Independent", "High", "Data Analysis", "Communication"],
                ["Data Analyst", "Ad-hoc analysis for stakeholders", 6, 15, "Variable", "Probabilistic", "Interactive", "Moderate", "Critical Thinking", "Data Analysis"],
                ["Data Analyst", "Stakeholder presentations", 4, 10, "Variable", "Judgment-heavy", "Collaborative", "Low", "Communication", "Stakeholder Mgmt"],
                ["Financial Analyst", "Monthly close reconciliation", 12, 30, "Repetitive", "Deterministic", "Independent", "High", "Data Analysis", "Process Automation"],
            ],
            "notes": "Required for: Work Design Lab, AI Opportunity Scan, Impact Simulator. Task Type: Repetitive/Variable. Logic: Deterministic/Probabilistic/Judgment-heavy. Interaction: Independent/Interactive/Collaborative. AI Impact: High/Moderate/Low. Time Spent % should sum to 100 per job.",
        },
        "skills": {
            "sheet_name": "Skills",
            "headers": ["Employee ID", "Name", "Skill", "Proficiency"],
            "samples": [
                ["EMP001", "Sarah Chen", "Data Analysis", 4],
                ["EMP001", "Sarah Chen", "Python/SQL", 3],
                ["EMP001", "Sarah Chen", "AI/ML Tools", 2],
                ["EMP001", "Sarah Chen", "Communication", 3],
                ["EMP002", "James Rodriguez", "Data Analysis", 4],
                ["EMP002", "James Rodriguez", "Leadership", 3],
                ["EMP003", "Maria Thompson", "Data Analysis", 3],
                ["EMP003", "Maria Thompson", "Critical Thinking", 3],
            ],
            "notes": "Required for: Skills Inventory, Gap Analysis, Adjacency Map. Proficiency scale: 1=Novice, 2=Developing, 3=Proficient, 4=Expert. One row per employee per skill.",
        },
        "org_structure": {
            "sheet_name": "Org Structure",
            "headers": ["Employee ID", "Name", "Job Title", "Function", "Career Level", "Manager ID", "Span of Control", "Cost Center", "Location"],
            "samples": [
                ["EMP001", "Sarah Chen", "Data Analyst", "Technology", "L3", "EMP010", 0, "CC100", "New York"],
                ["EMP010", "Rachel Anderson", "VP Technology", "Technology", "L5", "EMP020", 5, "CC100", "San Francisco"],
                ["EMP020", "Andrew Scott", "CEO", "Executive", "L6", "", 5, "CC000", "New York"],
            ],
            "notes": "Required for: Org Design Studio, Manager Capability. Span of Control = number of direct reports. Manager ID links to another employee's Employee ID.",
        },
        "ai_readiness": {
            "sheet_name": "AI Readiness",
            "headers": ["Employee ID", "Name", "Function", "AI Awareness", "Tool Adoption", "Data Literacy", "Change Openness", "AI Collaboration"],
            "samples": [
                ["EMP001", "Sarah Chen", "Technology", 4, 3, 4, 3, 3],
                ["EMP002", "James Rodriguez", "Technology", 4, 4, 4, 3, 4],
                ["EMP003", "Maria Thompson", "Finance", 3, 2, 3, 3, 2],
                ["EMP004", "David Kim", "Operations", 2, 1, 2, 2, 1],
                ["EMP010", "Rachel Anderson", "Technology", 4, 3, 3, 4, 3],
            ],
            "notes": "Required for: AI Readiness Assessment, Change Readiness. Score each dimension 1-5. 1=No awareness, 2=Basic, 3=Intermediate, 4=Advanced, 5=Expert.",
        },
        "manager_capability": {
            "sheet_name": "Manager Capability",
            "headers": ["Employee ID", "Name", "Function", "Direct Reports", "Leading Through Ambiguity", "AI Tool Fluency", "Coaching & Reskilling", "Communication of Change"],
            "samples": [
                ["EMP010", "Rachel Anderson", "Technology", 5, 4, 4, 3, 4],
                ["EMP011", "Thomas Harris", "Finance", 3, 3, 2, 3, 3],
                ["EMP012", "Catherine Clark", "Operations", 4, 2, 1, 2, 2],
                ["EMP013", "Steven Wright", "HR", 3, 3, 3, 4, 4],
            ],
            "notes": "Required for: Manager Capability Assessment, Manager Development. Only include people managers (those with direct reports). Score each dimension 1-5.",
        },
        "operating_model": {
            "sheet_name": "Operating Model",
            "headers": ["Function", "Capability", "Current Maturity", "Target Maturity", "Service Model", "Owner"],
            "samples": [
                ["Finance", "Accounting & Close", 3, 4, "Shared", "VP Finance"],
                ["Finance", "FP&A", 2, 4, "Embedded", "VP Finance"],
                ["Finance", "Treasury", 2, 3, "Shared", "VP Finance"],
                ["Technology", "Data & Analytics", 3, 5, "Platform", "VP Technology"],
                ["Technology", "Infrastructure", 3, 4, "Shared", "VP Technology"],
            ],
            "notes": "Required for: Operating Model Lab. Maturity: 1=Ad Hoc, 2=Repeatable, 3=Defined, 4=Managed, 5=Optimized. Service Model: Shared/Embedded/CoE/Outsourced/Platform.",
        },
        "change_management": {
            "sheet_name": "Change Management",
            "headers": ["Employee ID", "Name", "Function", "Change Readiness", "Impact Level", "Support Needed", "Preferred Communication", "Manager ID"],
            "samples": [
                ["EMP001", "Sarah Chen", "Technology", 4, "High", "Light", "Digital/Email", "EMP010"],
                ["EMP003", "Maria Thompson", "Finance", 3, "Moderate", "Moderate", "Town Hall", "EMP011"],
                ["EMP004", "David Kim", "Operations", 2, "High", "Intensive", "1:1 Meeting", "EMP012"],
                ["EMP009", "Jennifer Garcia", "Legal", 3, "Low", "Light", "Newsletter", "EMP014"],
            ],
            "notes": "Required for: Change Readiness, Change Planner. Change Readiness: 1-5. Impact Level: High/Moderate/Low. Support Needed: Intensive/Moderate/Light.",
        },
    }
    
    tmpl = templates.get(template_type)
    if not tmpl:
        return {"error": f"Unknown template: {template_type}. Available: {list(templates.keys())}"}
    
    ws.title = tmpl["sheet_name"]
    ws.append(tmpl["headers"])
    
    # Add sample rows
    samples = tmpl.get("samples", [tmpl.get("sample", [])])
    for sample in samples:
        ws.append(sample)
    
    # Style headers
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    header_fill = PatternFill(start_color="1A2340", end_color="1A2340", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    thin_border = Border(
        bottom=Side(style='thin', color='3B82F6'),
    )
    for cell in ws[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = thin_border
    
    # Style sample data rows
    sample_font = Font(color="888888", italic=True, size=10)
    for row in ws.iter_rows(min_row=2, max_row=len(samples)+1):
        for cell in row:
            cell.font = sample_font
    
    # Auto-width
    for col in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = max(max_len + 4, 14)
    
    # Add notes sheet if available
    if tmpl.get("notes"):
        notes_ws = wb.create_sheet("Instructions")
        notes_ws.append(["Instructions for " + tmpl["sheet_name"]])
        notes_ws.append([])
        notes_ws.append([tmpl["notes"]])
        notes_ws.append([])
        notes_ws.append(["Column Definitions:"])
        for i, h in enumerate(tmpl["headers"]):
            notes_ws.append([f"  {h}", "— Required" if i < 3 else "— Recommended"])
        notes_ws.append([])
        notes_ws.append(["Tips:"])
        notes_ws.append(["  • Delete the sample rows before uploading your data"])
        notes_ws.append(["  • Do not change column header names"])
        notes_ws.append(["  • Keep the sheet name as-is"])
        notes_ws.column_dimensions['A'].width = 50
        notes_ws.column_dimensions['B'].width = 20
        notes_ws['A1'].font = Font(bold=True, size=12, color="3B82F6")
    
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={template_type}_template.xlsx"}
    )


# ═══════════════════════════════════════════════════
# SCENARIO NARRATIVE STORAGE (in-memory, per model)
# ═══════════════════════════════════════════════════

_narrative_store: dict[str, str] = {}

@router.post("/export/narrative/{model_id}")
async def save_narrative(model_id: str, body: dict):
    """Store the scenario narrative for inclusion in exports."""
    _narrative_store[model_id] = str(body.get("narrative", ""))
    return {"status": "ok"}

@router.get("/export/narrative/{model_id}")
async def get_narrative(model_id: str):
    """Retrieve stored scenario narrative."""
    return {"narrative": _narrative_store.get(model_id, "")}


# ═══════════════════════════════════════════════════
# DOCX REPORT GENERATION
# ═══════════════════════════════════════════════════

def _amber_heading(doc, text, level=1):
    """Add a heading styled with warm amber branding."""
    from docx.shared import RGBColor, Pt
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0xC0, 0x70, 0x30)
    return h

def _add_styled_table(doc, headers, rows_data):
    """Add a table with amber-styled header row."""
    from docx.shared import RGBColor, Pt
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = 'Table Grid'
    # Style header row
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = 1  # center
        # Shade header cell amber
        shading = OxmlElement('w:shd')
        shading.set(qn('w:fill'), 'C07030')
        cell._tc.get_or_add_tcPr().append(shading)
    for row_data in rows_data:
        row = table.add_row()
        for i, val in enumerate(row_data):
            row.cells[i].text = str(val)
            for p in row.cells[i].paragraphs:
                for run in p.runs:
                    run.font.size = Pt(9)
    return table


@router.get("/export/docx/{model_id}")
async def export_docx(model_id: str):
    """Generate a styled Word document with transformation findings."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import io, datetime

    ds = store.datasets.get(model_id)
    if not ds:
        raise HTTPException(404, "Model not found")

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(10)

    wf = ds.get("workforce", pd.DataFrame())
    wd = ds.get("work_design", pd.DataFrame())
    jc = ds.get("job_catalog", pd.DataFrame())
    op = ds.get("operating_model", pd.DataFrame())
    ch = ds.get("change_management", pd.DataFrame())
    sk = ds.get("skills", pd.DataFrame())

    employees = len(wf) if not wf.empty else 0
    roles = int(get_series(wf, "Job Title").nunique()) if not wf.empty and "Job Title" in wf.columns else 0
    tasks = len(wd) if not wd.empty else 0
    functions = int(get_series(wf, "Function ID").nunique()) if not wf.empty and "Function ID" in wf.columns else 0

    # ── Title Page ──
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run('\n\n\n\nAI Transformation Report\n')
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0xC0, 0x70, 0x30)
    sub_run = title_para.add_run(f'\n{model_id}\n{datetime.datetime.now().strftime("%B %d, %Y")}\n\nConfidential — For Internal Use Only')
    sub_run.font.size = Pt(12)
    sub_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    doc.add_page_break()

    # ── Section 1: Executive Summary ──
    _amber_heading(doc, '1. Executive Summary')
    doc.add_paragraph(f'This report covers the AI transformation assessment for {employees:,} employees across {functions} functions, {roles} unique roles, and {tasks} mapped tasks. It provides a comprehensive view of the current workforce, AI impact analysis, job architecture, skills landscape, and recommended transformation roadmap.')
    p = doc.add_paragraph()
    p.add_run('Key Metrics: ').bold = True
    p.add_run(f'{employees:,} employees | {roles} roles | {tasks} tasks mapped | {functions} functions')

    # ── Section 2: Workforce Overview ──
    _amber_heading(doc, '2. Workforce Overview')
    if not wf.empty and "Function ID" in wf.columns:
        func_counts = safe_value_counts(get_series(wf, "Function ID"))
        _add_styled_table(doc, ['Function', 'Headcount', '% of Org'], [
            [str(f), str(c), f'{round(c / max(employees, 1) * 100)}%'] for f, c in func_counts.items()
        ])
    if not wf.empty and "Career Track" in wf.columns:
        doc.add_paragraph('')
        track_counts = safe_value_counts(get_series(wf, "Career Track"))
        for track, count in track_counts.items():
            doc.add_paragraph(f'• {track}: {count} ({round(count / max(employees, 1) * 100)}%)', style='List Bullet')

    # ── Section 3: AI Impact Analysis ──
    _amber_heading(doc, '3. AI Impact Analysis')
    if not wd.empty and "AI Impact" in wd.columns:
        ai_counts = safe_value_counts(get_series(wd, "AI Impact"))
        total_tasks = int(ai_counts.sum())
        _add_styled_table(doc, ['AI Impact Level', 'Tasks', '% of Total'], [
            [str(imp), str(int(cnt)), f'{round(int(cnt) / max(total_tasks, 1) * 100)}%'] for imp, cnt in ai_counts.items()
        ])
        high_pct = round(ai_counts.get("High", 0) / max(total_tasks, 1) * 100)
        doc.add_paragraph(f'\n{high_pct}% of mapped tasks have high AI automation potential, representing significant opportunity for workforce transformation.')
    else:
        doc.add_paragraph('No work design data available for AI impact analysis.')

    # ── Section 4: Job Architecture ──
    _amber_heading(doc, '4. Job Architecture')
    if not wf.empty and "Job Title" in wf.columns:
        titles = get_series(wf, "Job Title").value_counts().head(15)
        _add_styled_table(doc, ['Job Title', 'Headcount'], [
            [str(t), str(c)] for t, c in titles.items()
        ])
    if not wf.empty and "Job Family" in wf.columns:
        doc.add_paragraph('')
        fam_counts = safe_value_counts(get_series(wf, "Job Family"))
        doc.add_paragraph(f'{len(fam_counts)} job families identified across the organization.')

    # ── Section 5: Skills Landscape ──
    _amber_heading(doc, '5. Skills Landscape')
    if not sk.empty and "Skill" in sk.columns:
        skill_counts = safe_value_counts(get_series(sk, "Skill")).head(10)
        _add_styled_table(doc, ['Skill', 'Employees with Skill'], [
            [str(s), str(c)] for s, c in skill_counts.items()
        ])
    else:
        doc.add_paragraph('Skills data not yet uploaded. Upload skills inventory for gap analysis.')

    # ── Section 6: Change Readiness ──
    _amber_heading(doc, '6. Change Readiness')
    if not ch.empty and "Change Readiness" in ch.columns:
        avg_readiness = round(get_series(ch, "Change Readiness").astype(float).mean(), 1)
        doc.add_paragraph(f'Average change readiness score: {avg_readiness}/5')
        if "Impact Level" in ch.columns:
            impact_counts = safe_value_counts(get_series(ch, "Impact Level"))
            for level, count in impact_counts.items():
                doc.add_paragraph(f'• {level} impact: {count} employees', style='List Bullet')
    else:
        doc.add_paragraph('Change readiness data not yet assessed.')

    # ── Section 7: Recommendations ──
    _amber_heading(doc, '7. Recommendations')
    recs = [
        'Prioritize high-AI-impact roles for Phase 1 transformation — focus on repetitive, deterministic tasks first',
        'Build reskilling pathways for affected employees before automation deployment',
        'Deploy change champions at 1:5 ratio across impacted teams for maximum adoption',
        'Implement phased rollout: Wave 1 (quick wins, months 1-3), Wave 2 (core redesign, months 4-8), Wave 3 (optimization, months 9-12)',
        'Establish AI governance framework with clear decision rights before tool deployment',
        'Create dedicated measurement framework tracking productivity, adoption, and employee sentiment',
    ]
    for r in recs:
        doc.add_paragraph(f'• {r}', style='List Bullet')

    # ── Section 8: Transformation Roadmap ──
    _amber_heading(doc, '8. Transformation Roadmap')
    phases = [
        ('Assess (M1-3)', 'Workforce baseline, AI readiness scan, stakeholder mapping, risk assessment'),
        ('Design (M2-6)', 'Role redesign, job architecture, operating model, skill gap analysis'),
        ('Pilot (M5-8)', 'Wave 1 rollout, training delivery, tool deployment, feedback collection'),
        ('Scale (M7-12)', 'Enterprise rollout, process optimization, manager development, performance tracking'),
        ('Sustain (M11-15)', 'Continuous improvement, capability maturity, knowledge transfer, BAU transition'),
    ]
    _add_styled_table(doc, ['Phase', 'Key Activities'], phases)

    # ── Section 9: Scenario Narrative (if available) ──
    saved_narrative = _narrative_store.get(model_id, "")
    if saved_narrative:
        _amber_heading(doc, '9. Scenario Analysis')
        for para in saved_narrative.split("\n\n"):
            if para.strip():
                doc.add_paragraph(para.strip())

    # ── Section 10: Next Steps ──
    _amber_heading(doc, '10. Next Steps')
    steps = [
        'Complete role redesign for top 10 highest-impact roles',
        'Finalize BBBA dispositions and headcount waterfall',
        'Build detailed change management roadmap with stakeholder sign-off',
        'Present transformation business case to steering committee',
        'Launch pilot program with selected champion teams',
    ]
    for i, s in enumerate(steps, 1):
        doc.add_paragraph(f'{i}. {s}')

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={model_id}_Transformation_Report.docx"}
    )


# ═══════════════════════════════════════════════════
# POWERPOINT EXPORT — 12-slide deck
# ═══════════════════════════════════════════════════

@router.get("/export/pptx/{model_id}")
async def export_pptx(model_id: str):
    """Generate a 12-slide PowerPoint deck with warm amber branding."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import io, datetime

    ds = store.datasets.get(model_id)
    if not ds:
        raise HTTPException(404, "Model not found")

    wf = ds.get("workforce", pd.DataFrame())
    wd = ds.get("work_design", pd.DataFrame())
    ch = ds.get("change_management", pd.DataFrame())
    sk = ds.get("skills", pd.DataFrame())
    op = ds.get("operating_model", pd.DataFrame())

    employees = len(wf) if not wf.empty else 0
    roles = int(get_series(wf, "Job Title").nunique()) if not wf.empty and "Job Title" in wf.columns else 0
    tasks = len(wd) if not wd.empty else 0
    functions = int(get_series(wf, "Function ID").nunique()) if not wf.empty and "Function ID" in wf.columns else 0

    # AI impact counts
    high_ai = mod_ai = low_ai = 0
    if not wd.empty and "AI Impact" in wd.columns:
        ai_vc = safe_value_counts(get_series(wd, "AI Impact"))
        high_ai = ai_vc.get("High", 0)
        mod_ai = ai_vc.get("Moderate", 0)
        low_ai = ai_vc.get("Low", 0)
    total_ai = max(high_ai + mod_ai + low_ai, 1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    AMBER = RGBColor(0xC0, 0x70, 0x30)
    GOLD = RGBColor(0xE8, 0xC5, 0x47)
    DARK = RGBColor(0x1A, 0x23, 0x40)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x88, 0x88, 0x88)

    def _add_bg(slide, r=26, g=35, b=64):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def _add_bullet_slide(slide, items, top=2.2, font_size=16):
        for i, item in enumerate(items):
            _add_text_box(slide, 0.8, top + i * 0.55, 11, 0.5, f"→  {item}", font_size=font_size, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 1, 1.5, 11, 1.5, "AI Transformation", font_size=44, bold=True, color=AMBER, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 1, 3, 11, 1, "Assessment Report", font_size=32, color=GOLD, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 1, 4.5, 11, 0.5, f"{model_id}  |  {datetime.datetime.now().strftime('%B %Y')}  |  Confidential", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 2: Executive Summary
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Executive Summary", font_size=28, bold=True, color=AMBER)
    metrics = [
        (f"{employees:,}", "Employees"), (str(roles), "Roles"), (str(tasks), "Tasks Mapped"),
        (str(functions), "Functions"), (f"{round(high_ai/total_ai*100)}%", "High AI Impact"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = 0.5 + i * 2.4
        _add_text_box(s, x, 1.3, 2.2, 0.7, val, font_size=36, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
        _add_text_box(s, x, 2.0, 2.2, 0.4, label, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 3, 12, 2, f"This assessment covers {employees:,} employees across {functions} functions. {round(high_ai/total_ai*100)}% of {tasks} mapped tasks show high AI automation potential, representing significant transformation opportunity.", font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Slide 3: Workforce Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Workforce Overview", font_size=28, bold=True, color=AMBER)
    if not wf.empty and "Function ID" in wf.columns:
        func_counts = safe_value_counts(get_series(wf, "Function ID"))
        items = [f"{f}: {c:,} employees ({round(c/max(employees,1)*100)}%)" for f, c in list(func_counts.items())[:8]]
        _add_bullet_slide(s, items, top=1.3)

    # Slide 4: AI Impact Analysis
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "AI Impact Analysis", font_size=28, bold=True, color=AMBER)
    _add_text_box(s, 0.5, 1.3, 3.5, 0.7, f"{high_ai}", font_size=48, bold=True, color=RGBColor(0xEF, 0x44, 0x44), alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 2.0, 3.5, 0.4, "High Impact Tasks", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 4.5, 1.3, 3.5, 0.7, f"{mod_ai}", font_size=48, bold=True, color=RGBColor(0xF5, 0x9E, 0x0B), alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 4.5, 2.0, 3.5, 0.4, "Moderate Impact Tasks", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 8.5, 1.3, 3.5, 0.7, f"{low_ai}", font_size=48, bold=True, color=RGBColor(0x10, 0xB9, 0x81), alignment=PP_ALIGN.CENTER)
    _add_text_box(s, 8.5, 2.0, 3.5, 0.4, "Low Impact Tasks", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 5: Job Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Job Architecture", font_size=28, bold=True, color=AMBER)
    if not wf.empty and "Job Title" in wf.columns:
        top_roles = get_series(wf, "Job Title").value_counts().head(8)
        items = [f"{t}: {c} employees" for t, c in top_roles.items()]
        _add_bullet_slide(s, items, top=1.3)

    # Slide 6: Skills Landscape
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Skills Landscape", font_size=28, bold=True, color=AMBER)
    if not sk.empty and "Skill" in sk.columns:
        top_skills = safe_value_counts(get_series(sk, "Skill")).head(8)
        items = [f"{s}: {c} employees proficient" for s, c in top_skills.items()]
        _add_bullet_slide(s, items, top=1.3)
    else:
        _add_text_box(s, 0.5, 2.5, 12, 1, "Skills inventory not yet uploaded", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 7: Change Readiness
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Change Readiness", font_size=28, bold=True, color=AMBER)
    if not ch.empty and "Change Readiness" in ch.columns:
        avg_rd = round(get_series(ch, "Change Readiness").astype(float).mean(), 1)
        _add_text_box(s, 3, 1.5, 7, 1, f"{avg_rd}/5", font_size=72, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
        _add_text_box(s, 3, 2.8, 7, 0.5, "Average Change Readiness Score", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    else:
        _add_text_box(s, 0.5, 2.5, 12, 1, "Change readiness not yet assessed", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 8: Operating Model
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Operating Model", font_size=28, bold=True, color=AMBER)
    if not op.empty:
        items = [f"{row.get('Function', row.get('Function ID', ''))}: {row.get('Capability', '')}" for _, row in op.head(8).iterrows()]
        _add_bullet_slide(s, items, top=1.3)
    else:
        _add_text_box(s, 0.5, 2.5, 12, 1, "Operating model not yet defined", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 9: Recommendations
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Recommendations", font_size=28, bold=True, color=AMBER)
    recs = [
        "Prioritize high-AI-impact roles for Phase 1 transformation",
        "Build reskilling pathways before automation deployment",
        "Deploy change champions at 1:5 ratio for adoption",
        "Implement phased rollout across 3 waves over 12 months",
        "Establish AI governance framework with clear decision rights",
        "Create measurement framework for productivity and sentiment",
    ]
    _add_bullet_slide(s, recs, top=1.3)

    # Slide 10: Roadmap Timeline
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Transformation Roadmap", font_size=28, bold=True, color=AMBER)
    phases = [
        ("Assess", "M1-3", "Baseline, readiness scan, stakeholder map"),
        ("Design", "M2-6", "Role redesign, job architecture, skills gap"),
        ("Pilot", "M5-8", "Wave 1, training, tool deployment"),
        ("Scale", "M7-12", "Enterprise rollout, process optimization"),
        ("Sustain", "M11-15", "Continuous improvement, BAU transition"),
    ]
    for i, (phase, time, desc) in enumerate(phases):
        y = 1.3 + i * 1.0
        _add_text_box(s, 0.5, y, 2.5, 0.4, phase, font_size=20, bold=True, color=GOLD)
        _add_text_box(s, 3.0, y, 1.5, 0.4, time, font_size=14, color=AMBER)
        _add_text_box(s, 4.8, y, 8, 0.4, desc, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Slide 11: Scenario Narrative (if available)
    saved_narrative = _narrative_store.get(model_id, "")
    if saved_narrative:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(s)
        _add_text_box(s, 0.5, 0.3, 5, 0.6, "Scenario Analysis", font_size=28, bold=True, color=AMBER)
        paras = [p.strip() for p in saved_narrative.split("\n\n") if p.strip()]
        # Headline
        if paras:
            _add_text_box(s, 0.5, 1.2, 12, 0.8, paras[0], font_size=16, bold=True, color=GOLD)
        # Body (fit remaining paragraphs)
        for i, para in enumerate(paras[1:4]):  # max 3 body paragraphs to fit
            _add_text_box(s, 0.5, 2.2 + i * 1.5, 12, 1.4, para, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

    # Slide 12: Investment & ROI
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Investment & ROI", font_size=28, bold=True, color=AMBER)
    roi_items = [
        f"Workforce in scope: {employees:,} employees across {functions} functions",
        f"Tasks with automation potential: {high_ai + mod_ai} of {tasks} ({round((high_ai+mod_ai)/max(tasks,1)*100)}%)",
        "Estimated capacity freed: 15-25% of FTE hours in Phase 1",
        "Target ROI: 3-5x within 18 months of full deployment",
        "Reskilling investment: $2-5K per affected employee",
    ]
    _add_bullet_slide(s, roi_items, top=1.3)

    # Slide 12: Next Steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s)
    _add_text_box(s, 0.5, 0.3, 5, 0.6, "Next Steps", font_size=28, bold=True, color=AMBER)
    steps = [
        "Complete role redesign for top 10 highest-impact roles",
        "Finalize BBBA dispositions and headcount waterfall",
        "Build change management roadmap with stakeholder sign-off",
        "Present transformation business case to steering committee",
        "Launch pilot program with selected champion teams",
    ]
    _add_bullet_slide(s, steps, top=1.3)
    _add_text_box(s, 2, 5.5, 9, 0.5, "Thank you", font_size=24, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={model_id}_Transformation_Deck.pptx"}
    )


# ═══════════════════════════════════════════════════
# PDF EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════

@router.get("/export/pdf/{model_id}")
async def export_pdf(model_id: str):
    """Generate a one-page PDF executive summary with warm amber branding."""
    from fpdf import FPDF
    import io, datetime

    ds = store.datasets.get(model_id)
    if not ds:
        raise HTTPException(404, "Model not found")

    wf = ds.get("workforce", pd.DataFrame())
    wd = ds.get("work_design", pd.DataFrame())
    ch = ds.get("change_management", pd.DataFrame())

    employees = len(wf) if not wf.empty else 0
    roles = int(get_series(wf, "Job Title").nunique()) if not wf.empty and "Job Title" in wf.columns else 0
    tasks = len(wd) if not wd.empty else 0
    functions = int(get_series(wf, "Function ID").nunique()) if not wf.empty and "Function ID" in wf.columns else 0

    high_ai = mod_ai = low_ai = 0
    if not wd.empty and "AI Impact" in wd.columns:
        ai_vc = safe_value_counts(get_series(wd, "AI Impact"))
        high_ai = ai_vc.get("High", 0)
        mod_ai = ai_vc.get("Moderate", 0)
        low_ai = ai_vc.get("Low", 0)
    total_ai = max(high_ai + mod_ai + low_ai, 1)

    avg_readiness = 0
    if not ch.empty and "Change Readiness" in ch.columns:
        avg_readiness = round(get_series(ch, "Change Readiness").astype(float).mean(), 1)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=False)

    # Header bar
    pdf.set_fill_color(192, 112, 48)  # Amber
    pdf.rect(0, 0, 210, 22, 'F')
    pdf.set_font("Helvetica", "B", 18)
    pdf.set_text_color(255, 255, 255)
    pdf.set_xy(10, 4)
    pdf.cell(0, 14, "AI Transformation - Executive Summary", align="L")

    # Subtitle
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(150, 150, 150)
    pdf.set_xy(10, 24)
    pdf.cell(0, 5, f"{model_id}  |  {datetime.datetime.now().strftime('%B %d, %Y')}  |  Confidential")

    # KPI boxes
    kpis = [
        (f"{employees:,}", "Employees"),
        (str(roles), "Roles"),
        (str(tasks), "Tasks"),
        (f"{round(high_ai/total_ai*100)}%", "High AI"),
        (f"{avg_readiness}/5", "Readiness"),
    ]
    box_w = 35
    x_start = 12
    pdf.set_xy(x_start, 33)
    for i, (val, label) in enumerate(kpis):
        x = x_start + i * (box_w + 4)
        pdf.set_fill_color(248, 244, 239)
        pdf.rect(x, 33, box_w, 18, 'F')
        pdf.set_draw_color(232, 197, 71)
        pdf.rect(x, 33, box_w, 18, 'D')
        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(192, 112, 48)
        pdf.set_xy(x, 34)
        pdf.cell(box_w, 10, val, align="C")
        pdf.set_font("Helvetica", "", 7)
        pdf.set_text_color(120, 120, 120)
        pdf.set_xy(x, 43)
        pdf.cell(box_w, 5, label, align="C")

    # Section: Top Findings
    y = 58
    pdf.set_draw_color(232, 197, 71)
    pdf.line(10, y, 200, y)
    y += 3
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(192, 112, 48)
    pdf.set_xy(10, y)
    pdf.cell(90, 6, "TOP FINDINGS")
    pdf.set_xy(105, y)
    pdf.cell(90, 6, "KEY RECOMMENDATIONS")
    y += 8

    findings = [
        f"{employees:,} employees assessed across {functions} functions",
        f"{high_ai} tasks ({round(high_ai/total_ai*100)}%) have high automation potential",
        f"Average change readiness: {avg_readiness}/5",
    ]
    recs = [
        "Phase 1: Automate high-impact repetitive tasks",
        "Deploy change champions at 1:5 ratio",
        "Build reskilling pathways for affected roles",
    ]

    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(60, 60, 60)
    for i, f in enumerate(findings):
        pdf.set_xy(12, y + i * 7)
        pdf.cell(85, 6, f"  {f}")
    for i, r in enumerate(recs):
        pdf.set_xy(107, y + i * 7)
        pdf.cell(85, 6, f"  {r}")

    # Section: FTE Impact & ROI
    y2 = y + max(len(findings), len(recs)) * 7 + 6
    pdf.set_draw_color(232, 197, 71)
    pdf.line(10, y2, 200, y2)
    y2 += 3
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(192, 112, 48)
    pdf.set_xy(10, y2)
    pdf.cell(0, 6, "FTE IMPACT & ROI PROJECTION")
    y2 += 8
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(60, 60, 60)
    fte_items = [
        f"Estimated FTE capacity freed: {round(high_ai * 0.4 + mod_ai * 0.2)} FTE-equivalent hours/week",
        f"Target ROI: 3-5x within 18 months  |  Reskilling: $2-5K per employee",
        f"Net headcount impact: Redeploy {round(high_ai * 0.3)} roles, upskill {round(mod_ai * 0.5)}",
    ]
    for i, item in enumerate(fte_items):
        pdf.set_xy(12, y2 + i * 6)
        pdf.cell(0, 5, f"  {item}")

    # Section: Roadmap
    y3 = y2 + len(fte_items) * 6 + 6
    pdf.set_draw_color(232, 197, 71)
    pdf.line(10, y3, 200, y3)
    y3 += 3
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(192, 112, 48)
    pdf.set_xy(10, y3)
    pdf.cell(0, 6, "TRANSFORMATION ROADMAP")
    y3 += 8
    phases = [
        ("Assess", "M1-3", 192, 112, 48),
        ("Design", "M2-6", 232, 197, 71),
        ("Pilot", "M5-8", 192, 96, 48),
        ("Scale", "M7-12", 217, 119, 6),
        ("Sustain", "M11+", 184, 96, 42),
    ]
    bar_w = 34
    for i, (name, time, r, g, b) in enumerate(phases):
        x = 12 + i * (bar_w + 3)
        pdf.set_fill_color(r, g, b)
        pdf.rect(x, y3, bar_w, 12, 'F')
        pdf.set_font("Helvetica", "B", 9)
        pdf.set_text_color(255, 255, 255)
        pdf.set_xy(x, y3 + 1)
        pdf.cell(bar_w, 5, name, align="C")
        pdf.set_font("Helvetica", "", 7)
        pdf.set_xy(x, y3 + 6)
        pdf.cell(bar_w, 5, time, align="C")

    # Footer
    pdf.set_font("Helvetica", "", 7)
    pdf.set_text_color(180, 180, 180)
    pdf.set_xy(10, 282)
    pdf.cell(0, 5, f"Generated by AI Transformation Platform  |  {datetime.datetime.now().strftime('%Y-%m-%d')}  |  Confidential", align="C")

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={model_id}_Executive_Summary.pdf"}
    )


# ═══════════════════════════════════════════════════
# PER-MODULE DATA EXPORTS
# ═══════════════════════════════════════════════════

@router.get("/export/module/{model_id}/{module_name}")
async def export_module_data(model_id: str, module_name: str):
    """Export a specific module's data as Excel."""
    import io

    ds = store.datasets.get(model_id)
    if not ds:
        raise HTTPException(404, "Model not found")

    # Map module names to dataset keys
    module_map = {
        "snapshot": "workforce", "workforce": "workforce",
        "jobs": "job_catalog", "job_catalog": "job_catalog",
        "design": "work_design", "work_design": "work_design",
        "org": "org_design", "org_design": "org_design",
        "opmodel": "operating_model", "operating_model": "operating_model",
        "plan": "change_management", "change_management": "change_management",
        "market": "market_data", "market_data": "market_data",
    }

    dataset_key = module_map.get(module_name, module_name)
    df = ds.get(dataset_key, pd.DataFrame())

    if df is None or (isinstance(df, pd.DataFrame) and df.empty):
        raise HTTPException(404, f"No data for module: {module_name}")

    buf = io.BytesIO()
    buf.write(dataframe_to_excel_bytes(df, module_name[:31]))
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={module_name}_{model_id}.xlsx"}
    )

